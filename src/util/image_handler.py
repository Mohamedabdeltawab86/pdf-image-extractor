import os
from PIL import Image
import io
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import subprocess
import platform
import fitz
from PyQt5.QtCore import QThread, pyqtSignal
<<<<<<< HEAD
from PIL import ImageEnhance
from PIL import ImageOps
import tempfile
import sys
=======
>>>>>>> recover-features


def invert_image(image_bytes):
    """
    Invert image using Pillow's built-in operations.

    Args:
        image_bytes (bytes): Original image bytes

    Returns:
        bytes: Inverted image bytes
    """
    # Convert bytes to image
    img = Image.open(io.BytesIO(image_bytes))

    # Convert to RGB if not already
    if img.mode != "RGB":
        img = img.convert("RGB")

    # Invert the image using Pillow's built-in operation
    inverted_img = Image.eval(img, lambda x: 255 - x)

    # Save to bytes
    output_buffer = io.BytesIO()
    inverted_img.save(output_buffer, format="JPEG", quality=95)
    return output_buffer.getvalue()


def remove_black_background(image):
    """Remove the black background from an image."""
    if image.mode != "RGBA":
        image = image.convert("RGBA")
    data = image.getdata()

    new_data = []
    for item in data:
        if item[0] < 50 and item[1] < 50 and item[2] < 50:
            new_data.append((255, 255, 255, 0))
        else:
            new_data.append(item)

    image.putdata(new_data)
    return image


def save_image(image_bytes, output_dir, image_filename, should_invert=False):
    """
    Save image with optional inversion
    """
    try:
        # Convert bytes to image
        img = Image.open(io.BytesIO(image_bytes))

        # Convert to RGB if not already
        if img.mode != "RGB":
            img = img.convert("RGB")

        # Invert if requested
        if should_invert:
            img = remove_black_background(img)

        # Save to output
        base_name = os.path.splitext(image_filename)[0]
        output_filename = f"{base_name}.jpg"
        image_path = os.path.join(output_dir, output_filename)

        img.save(image_path, "JPEG", quality=95)
        print(f"Saved {'inverted' if should_invert else 'original'} {output_filename}")
        return True

    except Exception as e:
        print(f"Error saving image {image_filename}: {str(e)}")
        return False


def open_file(filepath):
    """Open a file with the default system application"""
    print(f"\nAttempting to open file: {filepath}")
    try:
        if platform.system() == "Darwin":  # macOS
            print("Using macOS open command")
            subprocess.call(("open", filepath))
        elif platform.system() == "Windows":  # Windows
            print("Using Windows start command")
            os.startfile(filepath)
        else:  # Linux variants
            print("Using Linux xdg-open command")
            subprocess.call(("xdg-open", filepath))
        print("File opened successfully")
    except Exception as e:
        print(f"Error opening file: {str(e)}", file=sys.stderr)


def extract_to_ppt(images, output_dir, output_name, should_invert=False):
    """Extract images to PowerPoint presentation"""
    try:
        prs = Presentation()

        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Set default slide background to black
        for layout in prs.slide_layouts:
            background = layout.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

        for i, image_bytes in enumerate(images):
            # Add a slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

            # Process image
            img = Image.open(io.BytesIO(image_bytes))

            # Convert to RGB if needed
            if img.mode != "RGB":
                img = img.convert("RGB")

            # Invert if requested
            if should_invert:
                img = remove_black_background(img)

            # Save temporary file
            temp_path = os.path.join(output_dir, f"temp_{i}.png")
            img.save(temp_path, "PNG")

            # Calculate image dimensions to fit slide
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height

            # Set maximum dimensions (leaving margins)
            max_width = Inches(12)  # 1.333 inch margin total
            max_height = Inches(6.5)  # 1 inch margin total

            # Calculate dimensions maintaining aspect ratio
            if aspect_ratio > max_width / max_height:
                width = max_width
                height = width / aspect_ratio
            else:
                height = max_height
                width = height * aspect_ratio

            # Center the image on slide
            left = (prs.slide_width - width) / 2
            top = (prs.slide_height - height) / 2

            # Add to slide
            pic = slide.shapes.add_picture(temp_path, left, top, width, height)

            # Remove temp file
            os.remove(temp_path)

        # Save presentation with new naming
        ppt_path = os.path.join(output_dir, f"{output_name}.pptx")
        prs.save(ppt_path)

        # Open the file
        open_file(ppt_path)
        return len(images)

    except Exception as e:
        print(f"Error creating PowerPoint: {str(e)}")
        raise e


class ImageExtractionThread(QThread):
<<<<<<< HEAD
    progress = pyqtSignal(object)
    finished = pyqtSignal(tuple)

    def __init__(self, pdf_path, output_dir, start_page, end_page, options):
        super().__init__()
        self.pdf_path = pdf_path
        self.output_dir = output_dir
        self.start_page = start_page - 1
        self.end_page = end_page
        self.options = options
        self._is_running = True

        # Print initialization info
        print("\n=== Extraction Thread Initialized ===")
        print(f"PDF Path: {pdf_path}")
        print(f"Output Directory: {output_dir}")
        print(f"Page Range: {start_page} to {end_page}")
        print("=====================================\n")

        # Ensure output directory exists
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            print(f"Created output directory: {output_dir}")

    def open_file(self, filepath):
        """Open a file with the default system application"""
        print(f"\nAttempting to open file: {filepath}")
        try:
            if platform.system() == "Darwin":  # macOS
                print("Using macOS open command")
                subprocess.call(("open", filepath))
            elif platform.system() == "Windows":  # Windows
                print("Using Windows start command")
                os.startfile(filepath)
            else:  # Linux variants
                print("Using Linux xdg-open command")
                subprocess.call(("xdg-open", filepath))
            print("File opened successfully")
        except Exception as e:
            print(f"Error opening file: {str(e)}", file=sys.stderr)
=======
    progress = pyqtSignal(int)
    finished = pyqtSignal(tuple)

    def __init__(self, pdf_path, output_dir, start_page=1, end_page=None, options=None):
        super().__init__()
        self.pdf_path = pdf_path
        self.output_dir = output_dir
        self.start_page = start_page
        self.end_page = end_page
        self.options = options or {}
        self._is_running = True

    def remove_black_background(self, image_bytes):
        """Remove black background more accurately from an image."""
        img = Image.open(io.BytesIO(image_bytes)).convert("RGBA")
        data = img.getdata()

        new_data = []
        for item in data:
            # Check if the pixel is predominantly black (all channels below 50)
            if item[0] < 50 and item[1] < 50 and item[2] < 50:
                new_data.append((255, 255, 255, 0))  # Fully transparent
            else:
                # Keep original pixel with full opacity
                new_data.append((item[0], item[1], item[2], 255))

        img.putdata(new_data)
        buffer = io.BytesIO()
        img.save(buffer, format="PNG")
        return buffer.getvalue()

    def merge_images(self, original_bytes, annotation_bytes):
        """Merge original image with annotation while preserving transparency."""
        original_image = Image.open(io.BytesIO(original_bytes))
        if original_image.mode != "RGBA":
            original_image = original_image.convert("RGBA")

        annotation_image = Image.open(io.BytesIO(annotation_bytes))
        if annotation_image.mode != "RGBA":
            annotation_image = annotation_image.convert("RGBA")

        # Ensure both images are the same size
        if annotation_image.size != original_image.size:
            annotation_image = annotation_image.resize(
                original_image.size, Image.Resampling.LANCZOS
            )

        # Create a new blank image with alpha channel
        merged_image = Image.new("RGBA", original_image.size, (0, 0, 0, 0))

        # Paste original image first
        merged_image.paste(original_image, (0, 0))

        # Composite annotation on top with alpha channel
        merged_image = Image.alpha_composite(merged_image, annotation_image)

        buffer = io.BytesIO()
        merged_image.save(buffer, format="PNG")
        return buffer.getvalue()

    def process_page_images(self, page, doc):
        image_groups = {}
        image_list = page.get_images(full=True)

        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]

            # Get image rectangle for positioning
            image_rect = page.get_image_bbox(img)
            rect_key = (round(image_rect.x0, 2), round(image_rect.y0, 2))

            if rect_key not in image_groups:
                image_groups[rect_key] = []
            image_groups[rect_key].append((image_bytes, base_image["ext"]))

        processed_images = []
        for rect_key, images in image_groups.items():
            if len(images) == 2:
                # Identify annotation layer
                img1_bytes, _ = images[0]
                img2_bytes, _ = images[1]

                img1 = Image.open(io.BytesIO(img1_bytes))
                img2 = Image.open(io.BytesIO(img2_bytes))

                # Check which image is likely the annotation
                is_img1_annotation = self.is_annotation_layer(img1)
                is_img2_annotation = self.is_annotation_layer(img2)

                if is_img2_annotation:
                    original_bytes, annotation_bytes = img1_bytes, img2_bytes
                else:
                    original_bytes, annotation_bytes = img2_bytes, img1_bytes

                # Process and merge
                cleaned_annotation = self.remove_black_background(annotation_bytes)
                merged_image = self.merge_images(original_bytes, cleaned_annotation)

                # Get caption if exists
                caption = self.get_caption(page, image_rect)
                processed_images.append((merged_image, caption))
            elif self.options.get("include_non_annotated", True):
                # Single image without annotation
                processed_images.append((images[0][0], ""))

        return processed_images

    def is_annotation_layer(self, img):
        """Determine if an image is likely an annotation layer."""
        if img.mode == "RGBA":
            return True

        # Convert to RGB if needed
        if img.mode != "RGB":
            img = img.convert("RGB")

        # Sample pixels to check if it's mostly black/white
        pixels = img.getdata()
        black_white_count = 0
        sample_size = min(1000, len(pixels))

        for i in range(0, len(pixels), len(pixels) // sample_size):
            p = pixels[i]
            # Check if pixel is very dark or very light
            if (p[0] < 30 and p[1] < 30 and p[2] < 30) or (
                p[0] > 225 and p[1] > 225 and p[2] > 225
            ):
                black_white_count += 1

        return (black_white_count / sample_size) > 0.8

    def get_caption(self, page, image_rect):
        """Extract caption text below the image."""
        caption_rect = fitz.Rect(
            image_rect.x0, image_rect.y1, image_rect.x1, page.rect.height
        )
        caption_text = page.get_textbox(caption_rect)
        return caption_text.strip()

    def extract_to_ppt(self, processed_images, output_path, include_non_annotated=True):
        """Extract images to PowerPoint file."""
        temp_files = []  # Keep track of temporary files
        try:
            prs = Presentation()

            # Set slide dimensions (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Set default slide background to black
            for layout in prs.slide_layouts:
                background = layout.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 0, 0)

            for i, (image_bytes, caption) in enumerate(processed_images):
                # Create unique temporary file name
                temp_path = os.path.join(
                    self.output_dir, f"temp_image_{i}_{os.getpid()}.png"
                )
                temp_files.append(temp_path)  # Add to cleanup list

                # Save image to temporary file
                with open(temp_path, "wb") as f:
                    f.write(image_bytes)

                # Get image dimensions before adding to slide
                with Image.open(temp_path) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height

                # Add to PowerPoint
                slide = prs.slides.add_slide(prs.slide_layouts[5])

                # Set maximum dimensions (leaving margins)
                max_width = Inches(12)  # 1.333 inch margin total
                max_height = Inches(6.75)  # 0.75 inch margin total

                # Calculate dimensions maintaining aspect ratio
                if aspect_ratio > max_width / max_height:
                    width = max_width
                    height = width / aspect_ratio
                else:
                    height = max_height
                    width = height * aspect_ratio

                # Center the image on slide
                left = (prs.slide_width - width) / 2
                top = (prs.slide_height - height) / 2

                # Add picture with calculated dimensions
                slide.shapes.add_picture(
                    temp_path, left, top, width=width, height=height
                )

                # Add caption if exists
                if caption:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = caption

            # Save PowerPoint file
            prs.save(output_path)

            # Clean up all temporary files
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception as e:
                    print(
                        f"Warning: Could not remove temporary file {temp_file}: {str(e)}"
                    )

            # Open the PowerPoint file
            self.open_file(output_path)

            return True

        except Exception as e:
            print(f"Error creating PowerPoint: {str(e)}")
            # Attempt to clean up on error
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except:
                    pass
            return False
>>>>>>> recover-features

    def run(self):
        doc = None
        try:
<<<<<<< HEAD
            print("\n=== Starting PDF Processing ===")
            doc = fitz.open(self.pdf_path)
            prs = Presentation()

            print(f"Processing PDF: {self.pdf_path}")
            print(f"Total pages in range: {self.end_page - self.start_page}")

            total_pages = self.end_page - self.start_page
            processed_pages = 0
            images_added = 0

            with tempfile.TemporaryDirectory() as temp_dir:
                print(f"Created temp directory: {temp_dir}")

                for page_num in range(self.start_page, self.end_page):
                    if not self._is_running:
                        print("Processing stopped by user")
                        break

                    processed_pages += 1
                    print(f"\nProcessing page {page_num + 1}")
                    self.progress.emit(
                        (
                            processed_pages,
                            total_pages,
                            f"معالجة الصفحة {page_num + 1} من {self.end_page}",
                        )
                    )

                    for img in doc.get_page_images(page_num):
                        try:
                            xref = img[0]
                            image = doc.extract_image(xref)

                            if not image:
                                print(
                                    f"No image data for image {xref} on page {page_num + 1}"
                                )
                                continue

                            temp_path = os.path.join(
                                temp_dir, f"temp_{page_num}_{xref}.png"
                            )
                            print(f"Saving temp image: {temp_path}")

                            with open(temp_path, "wb") as tmp:
                                tmp.write(image["image"])

                            slide = prs.slides.add_slide(prs.slide_layouts[6])
                            slide.shapes.add_picture(
                                temp_path, 0, 0, prs.slide_width, prs.slide_height
                            )
                            images_added += 1
                            print(f"Added image {images_added} to PowerPoint")

                        except Exception as e:
                            print(
                                f"Error processing image on page {page_num + 1}: {str(e)}",
                                file=sys.stderr,
                            )
                            continue

                if self._is_running and images_added > 0:
                    output_path = os.path.join(self.output_dir, "extracted_images.pptx")
                    print(f"\nSaving PowerPoint to: {output_path}")

                    prs.save(output_path)
                    print("PowerPoint file saved successfully")

                    self.msleep(500)

                    if os.path.exists(output_path):
                        file_size = os.path.getsize(output_path)
                        print(f"File size: {file_size} bytes")
                        if file_size > 0:
                            print("Opening PowerPoint file...")
                            self.open_file(output_path)
                            self.finished.emit(
                                (
                                    True,
                                    f"تم حفظ {images_added} صورة في العرض التقديمي بنجاح",
                                    images_added,
                                )
                            )
                        else:
                            print("Error: File is empty")
                            self.finished.emit((False, "تم إنشاء ملف فارغ", 0))
                    else:
                        print("Error: File not found after saving")
                        self.finished.emit(
                            (False, "لم يتم العثور على الملف بعد الحفظ", 0)
                        )
                else:
                    print("No images were processed")
                    self.finished.emit((False, "لم يتم العثور على صور للاستخراج", 0))

        except Exception as e:
            print(f"Error in run(): {str(e)}", file=sys.stderr)
            self.finished.emit((False, str(e), 0))
        finally:
            if doc:
                doc.close()
                print("PDF document closed")
            print("=== Processing Complete ===\n")

    def save_images(self, images):
        """Save images as separate files"""
        for i, img in enumerate(images, 1):
            try:
                filename = f"image_{i}.{self.options.get('format', 'PNG').lower()}"
                filepath = os.path.join(self.output_dir, filename)
                img.save(filepath)
            except Exception as e:
                print(f"Error saving image {i}: {str(e)}")

    def enhance_image(self, img):
        """Enhance image quality"""
        if self.options.get("enhance"):
            img = img.convert("RGB")
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(1.2)
            enhancer = ImageEnhance.Sharpness(img)
            img = enhancer.enhance(1.1)
        return img

    def remove_background(self, img):
        """Remove image background"""
        if self.options.get("remove_bg"):
            try:
                img = img.convert("RGBA")
                data = img.getdata()
                new_data = []
                for item in data:
                    if item[0] > 240 and item[1] > 240 and item[2] > 240:
                        new_data.append((255, 255, 255, 0))
                    else:
                        new_data.append(item)
                img.putdata(new_data)
            except Exception:
                pass
        return img

    @staticmethod
    def open_file(path):
        """Open a file or directory using the default application"""
        try:
            if os.path.exists(path):
                if os.name == "nt":  # Windows
                    os.startfile(path)
                elif os.name == "posix":  # macOS and Linux
                    if os.path.isfile(path):
                        subprocess.Popen(["xdg-open", path])
                    else:
                        subprocess.Popen(["xdg-open", path])
=======
            doc = fitz.open(self.pdf_path)
            doc_page_count = len(doc)

            if self.options.get("preview_only"):
                images = []
                start = self.start_page - 1
                end = self.end_page if self.end_page else doc_page_count

                for page_num in range(start, end):
                    if not self._is_running:
                        break

                    page = doc[page_num]
                    processed_images = self.process_page_images(page, doc)
                    images.extend([img for img, _ in processed_images])

                self.finished.emit((True, "تم استخراج الصور للمعاينة", images))
                return

            # For actual extraction
            start = self.start_page - 1
            end = self.end_page if self.end_page else doc_page_count

            all_processed_images = []
            for page_num in range(start, end):
                if not self._is_running:
                    break

                page = doc[page_num]
                processed_images = self.process_page_images(page, doc)
                all_processed_images.extend(processed_images)

            if all_processed_images:
                output_name = os.path.splitext(os.path.basename(self.pdf_path))[0]
                output_path = os.path.join(self.output_dir, f"{output_name}.pptx")

                success = self.extract_to_ppt(
                    all_processed_images,
                    output_path,
                    include_non_annotated=self.options.get(
                        "include_non_annotated", True
                    ),
                )

                if success:
                    self.finished.emit(
                        (
                            True,
                            f"تم حفظ {len(all_processed_images)} صورة",
                            len(all_processed_images),
                        )
                    )
                else:
                    self.finished.emit((False, "فشل في حفظ الملف", 0))
            else:
                self.finished.emit((False, "لم يتم العثور على صور", 0))

        except Exception as e:
            print(f"Error during extraction: {str(e)}")
            self.finished.emit((False, "حدث خطأ أثناء المعالجة", 0))
        finally:
            if doc:
                doc.close()

    def open_file(self, filepath):
        """Open a file with the default system application"""
        try:
            if platform.system() == "Windows":
                os.startfile(filepath)
            elif platform.system() == "Darwin":  # macOS
                subprocess.call(("open", filepath))
            else:  # Linux
                subprocess.call(("xdg-open", filepath))
>>>>>>> recover-features
        except Exception as e:
            print(f"Error opening file: {str(e)}")
