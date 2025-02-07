# src/modules/image_extractor/ppt_extractor.py

from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image
import os  # Import the os module
from datetime import datetime


def extract_to_ppt(images, output_dir, should_invert, base_filename):
    """
    Creates a PowerPoint presentation from a list of image byte strings.

    Args:
        images: A list of image byte strings.
        output_dir: The directory to save the PPTX file.
        should_invert:  Whether to invert the colors of the images.
        base_filename: The original PDF filename (without extension).
    """
    try:
        prs = Presentation()
        prs.slide_width = Inches(16)  # Standard widescreen size
        prs.slide_height = Inches(9)

        for image_bytes in images:
            try:
                img = Image.open(io.BytesIO(image_bytes))

                # Invert colors if requested
                if should_invert:
                    if img.mode != "RGB":
                        img = img.convert("RGB")  # Convert to RGB for inversion
                    img = Image.eval(img, lambda x: 255 - x)

                # Convert PIL Image to bytes for PPTX
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format='PNG')  # Save as PNG to preserve transparency
                img_byte_arr = img_byte_arr.getvalue()

                # Create a blank slide
                blank_slide_layout = prs.slide_layouts[6]  # 6 is the index for a blank slide
                slide = prs.slides.add_slide(blank_slide_layout)

                # Add the image to the slide
                left = top = Inches(0)
                pic = slide.shapes.add_picture(io.BytesIO(img_byte_arr), left, top, width=prs.slide_width, height=prs.slide_height)

            except Exception as e:
                print(f"Error processing image for PPT: {e}")
                continue  # Skip to the next image

        # Create the output filename with date and time
        now = datetime.now()
        date_time_str = now.strftime("%Y-%m-%d_%H-%M-%S")
        output_filename = f"{base_filename}_{date_time_str}.pptx"
        output_path = os.path.join(output_dir, output_filename)

        prs.save(output_path)
        print(f"PowerPoint saved to: {output_path}")

    except Exception as e:
        print(f"Error creating PowerPoint: {e}")